"""
Build the SQL on HPC lecture slide deck.
Morehouse-branded PPTX. Run: python3 build_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from datetime import date

# Morehouse brand colors
MAROON = RGBColor(0x84, 0x00, 0x28)
GOLD = RGBColor(0xC1, 0xA2, 0x31)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COOL_GRAY = RGBColor(0xA7, 0xA8, 0xAA)
LIGHT_GRAY = RGBColor(0xEE, 0xEE, 0xEE)

HEAD_FONT = "Arial"
BODY_FONT = "Arial"
MONO_FONT = "Courier New"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK = prs.slide_layouts[6]


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, font=BODY_FONT, size=18,
             bold=False, color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def add_paragraph(tf, text, font=BODY_FONT, size=18, bold=False, color=BLACK,
                  align=PP_ALIGN.LEFT, space_before=6):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def add_bullets(slide, left, top, width, height, items, font=BODY_FONT, size=20,
                color=BLACK, bullet_color=GOLD):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        # bullet marker
        run_b = p.add_run()
        run_b.text = "■  "
        run_b.font.name = font
        run_b.font.size = Pt(size)
        run_b.font.color.rgb = bullet_color
        run_b.font.bold = True
        # text
        run = p.add_run()
        run.text = item
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tf


def header_bar(slide, title_text, color=MAROON):
    """Maroon header bar with white title text."""
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(1.0), color)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12.5), Inches(0.7),
             title_text, font=HEAD_FONT, size=30, bold=True, color=WHITE)


def footer(slide, page_num=None, total=None):
    add_rect(slide, Emu(0), Inches(7.2), SLIDE_W, Inches(0.3), MAROON)
    add_text(slide, Inches(0.5), Inches(7.22), Inches(8), Inches(0.26),
             "SQL on HPC  |  Morehouse Supercomputing Facility",
             font=BODY_FONT, size=10, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    if page_num and total:
        add_text(slide, Inches(11.5), Inches(7.22), Inches(1.5), Inches(0.26),
                 f"{page_num} / {total}", font=BODY_FONT, size=10, color=WHITE,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


TOTAL_SLIDES = 24
slide_num = 0

# =============================================================================
# Slide 1: Title
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, MAROON)
# gold accent bar
add_rect(s, Inches(0.5), Inches(3.0), Inches(2), Inches(0.1), GOLD)
add_text(s, Inches(0.5), Inches(1.5), Inches(12), Inches(1.5),
         "SQL ON HPC", font=HEAD_FONT, size=72, bold=True, color=WHITE)
add_text(s, Inches(0.5), Inches(2.4), Inches(12), Inches(0.6),
         "Querying Large Databases on a Supercomputer",
         font=BODY_FONT, size=24, color=GOLD)
add_text(s, Inches(0.5), Inches(4.5), Inches(12), Inches(0.5),
         "Ashley Scruse, Ph.D.",
         font=HEAD_FONT, size=24, bold=True, color=WHITE)
add_text(s, Inches(0.5), Inches(5.0), Inches(12), Inches(0.4),
         "Deputy Director, Morehouse Supercomputing Facility",
         font=BODY_FONT, size=16, color=WHITE)
add_text(s, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
         "Guest Lecture  |  Database Systems  |  April 2026",
         font=BODY_FONT, size=14, color=COOL_GRAY)

# =============================================================================
# Slide 2: Who I am
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "WHO I AM")
add_text(s, Inches(0.5), Inches(1.5), Inches(12), Inches(0.6),
         "Ashley Scruse, Ph.D.", font=HEAD_FONT, size=32, bold=True, color=MAROON)
add_text(s, Inches(0.5), Inches(2.1), Inches(12), Inches(0.4),
         "Deputy Director, Morehouse Supercomputing Facility (MSCF)",
         font=BODY_FONT, size=18, color=BLACK)
add_bullets(s, Inches(0.5), Inches(3.0), Inches(12), Inches(3.5), [
    "I'm not a SQL expert. I'm here to bring the HPC piece.",
    "Bioinformatics Ph.D. from UGA via the Clark Atlanta HBCU pipeline.",
    "MSCF runs Morehouse's relationship with the Vista supercomputer at TACC.",
    "Today I'm wearing the data analyst hat with you.",
], size=20)
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 3: Roadmap
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "ROADMAP")
sections = [
    ("1", "HPC Intro", "What's a supercomputer, why a data analyst cares"),
    ("2", "Today's Frame", "You're an analyst, here's the data, here's the thinking"),
    ("3", "Live Demo", "We open Vista together and run real queries"),
    ("4", "The Investigation", "6 queries, 5 phases, build the case"),
    ("5", "Recommendation", "What does the data support?"),
    ("6", "Homework + Q&A", "You take a different question to the data"),
]
y = Inches(1.5)
for num, title, desc in sections:
    # number block
    add_rect(s, Inches(0.5), y, Inches(0.8), Inches(0.8), MAROON)
    add_text(s, Inches(0.5), y, Inches(0.8), Inches(0.8), num,
             font=HEAD_FONT, size=32, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.6), y + Inches(0.05), Inches(11), Inches(0.4),
             title, font=HEAD_FONT, size=22, bold=True, color=MAROON)
    add_text(s, Inches(1.6), y + Inches(0.45), Inches(11), Inches(0.35),
             desc, font=BODY_FONT, size=16, color=BLACK)
    y += Inches(0.95)
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 5: What is supercomputing?
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "WHAT IS SUPERCOMPUTING?")
add_text(s, Inches(0.5), Inches(1.4), Inches(12), Inches(0.6),
         "A bunch of really powerful computers wired together to act as one.",
         font=BODY_FONT, size=22, bold=True, color=MAROON)
add_bullets(s, Inches(0.5), Inches(2.4), Inches(12), Inches(4), [
    "Each computer in the cluster is called a node.",
    "Nodes have way more RAM and CPU than your laptop.",
    "You access the system over SSH and request resources to run your work.",
    "The system is shared among hundreds of users at the same time.",
], size=20)
add_text(s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
         "Vista (TACC): 500+ nodes, hundreds of GB RAM each, petabytes of storage.",
         font=BODY_FONT, size=16, color=COOL_GRAY)
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 6: Thousands of Computers Working Together (parallel workers table)
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "THOUSANDS OF COMPUTERS WORKING TOGETHER ON ONE PROBLEM")
add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.4),
         "100,000 math operations  |  5 seconds per operation",
         font=BODY_FONT, size=18, color=COOL_GRAY)

# Comparison table
table_top = Inches(1.9)
hr = Inches(0.65)
cw1 = Inches(2.6)
cw2 = Inches(3.2)
cw3 = Inches(3.5)
cw4 = Inches(3.0)
x1, x2, x3, x4 = Inches(0.5), Inches(3.1), Inches(6.3), Inches(9.8)

# Header row
add_rect(s, x1, table_top, cw1, hr, MAROON)
add_rect(s, x2, table_top, cw2, hr, MAROON)
add_rect(s, x3, table_top, cw3, hr, MAROON)
add_rect(s, x4, table_top, cw4, hr, GOLD)
add_text(s, x1, table_top, cw1, hr, "", font=HEAD_FONT, size=14, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x2 + Inches(0.05), table_top, cw2, hr, "1 Person\n(Your Laptop)",
         font=HEAD_FONT, size=12, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x3 + Inches(0.05), table_top, cw3, hr, "1,000 People\n(Small Cluster)",
         font=HEAD_FONT, size=12, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x4 + Inches(0.05), table_top, cw4, hr, "25,000 People\n(Supercomputer)",
         font=HEAD_FONT, size=12, bold=True, color=MAROON,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

rows = [
    ("Workers", "1", "1,000", "25,000"),
    ("Operations Each", "100,000", "100", "4"),
    ("Time to Complete", "500,000 sec", "500 sec", "20 sec"),
    ("In Minutes", "~8,333 min", "~8.3 min", "~0.3 min"),
    ("In Real Time", "~5.8 days", "~8 minutes", "20 seconds"),
]
y = table_top + hr
for i, (label, c1, c2, c3) in enumerate(rows):
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    is_last = (i == len(rows) - 1)
    label_color = MAROON
    add_rect(s, x1, y, cw1, hr, bg)
    add_rect(s, x2, y, cw2, hr, bg)
    add_rect(s, x3, y, cw3, hr, bg)
    add_rect(s, x4, y, cw4, hr, bg)
    add_text(s, x1 + Inches(0.1), y, cw1, hr, label,
             font=HEAD_FONT, size=14, bold=True, color=label_color,
             anchor=MSO_ANCHOR.MIDDLE)
    val_size = 16 if is_last else 14
    val_bold = is_last
    add_text(s, x2, y, cw2, hr, c1, font=BODY_FONT, size=val_size, bold=val_bold,
             color=BLACK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x3, y, cw3, hr, c2, font=BODY_FONT, size=val_size, bold=val_bold,
             color=BLACK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x4, y, cw4, hr, c3, font=BODY_FONT, size=val_size, bold=val_bold,
             color=MAROON, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += hr

add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
         "Same work. Same operations. The only difference is how many work in parallel.",
         font=HEAD_FONT, size=18, bold=True, color=MAROON, align=PP_ALIGN.CENTER)
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 7: Why should a data analyst care?
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "WHY SHOULD A DATA ANALYST CARE?")
# Comparison table: laptop vs HPC
table_top = Inches(1.5)
col1_x = Inches(0.5)
col2_x = Inches(4.7)
col3_x = Inches(9.0)
col_w = Inches(4.0)
row_h = Inches(0.7)

# Header row
add_rect(s, col1_x, table_top, col_w, row_h, COOL_GRAY)
add_rect(s, col2_x, table_top, col_w, row_h, MAROON)
add_rect(s, col3_x, table_top, col_w, row_h, MAROON)
add_text(s, col1_x, table_top, col_w, row_h, "What you need", font=HEAD_FONT,
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, col2_x, table_top, col_w, row_h, "Your laptop", font=HEAD_FONT,
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, col3_x, table_top, col_w, row_h, "HPC compute node", font=HEAD_FONT,
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

rows = [
    ("RAM", "8-16 GB", "128-223 GB"),
    ("Disk", "Limited SSD", "1 TB+ on $WORK"),
    ("38M-row scan", "Slow, swaps to disk", "Fast, fits in memory"),
    ("Many queries at once", "Bogs down", "Plenty of headroom"),
    ("Same env for the whole class", "Nope", "Yes"),
]
y = table_top + row_h
for i, (label, lap, hpc) in enumerate(rows):
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(s, col1_x, y, col_w, row_h, bg)
    add_rect(s, col2_x, y, col_w, row_h, bg)
    add_rect(s, col3_x, y, col_w, row_h, bg)
    add_text(s, col1_x + Inches(0.1), y, col_w, row_h, label,
             font=HEAD_FONT, size=14, bold=True, color=MAROON, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, col2_x + Inches(0.1), y, col_w, row_h, lap,
             font=BODY_FONT, size=14, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, col3_x + Inches(0.1), y, col_w, row_h, hpc,
             font=BODY_FONT, size=14, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
    y += row_h
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 7: Shared system
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "A SUPERCOMPUTER IS A SHARED SYSTEM")
add_text(s, Inches(0.5), Inches(1.5), Inches(12), Inches(0.7),
         "Hundreds of users on the same machine at the same time.",
         font=BODY_FONT, size=24, bold=True, color=MAROON)
add_bullets(s, Inches(0.5), Inches(2.6), Inches(12), Inches(4), [
    "This is why you can't run heavy queries on the login node. You'd slow down the whole system for everyone else.",
    "This is also why the file-sharing model works. One person can host a 6.3 GB database, set permissions, and a whole class can read it without copying it 35 times.",
    "Two sides of the same coin: shared resources require etiquette AND enable collaboration.",
], size=20)
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 8: Login vs compute nodes
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "LOGIN NODES VS COMPUTE NODES")
# two columns
col_y = Inches(1.5)
col_h = Inches(5)
# Login node card
add_rect(s, Inches(0.5), col_y, Inches(6), col_h, LIGHT_GRAY)
add_rect(s, Inches(0.5), col_y, Inches(6), Inches(0.6), COOL_GRAY)
add_text(s, Inches(0.5), col_y, Inches(6), Inches(0.6), "LOGIN NODE",
         font=HEAD_FONT, size=20, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(0.7), col_y + Inches(0.8), Inches(5.6), Inches(4), [
    "The first thing you land on when you SSH in",
    "Shared by everyone simultaneously",
    "Good for: editing files, copying, git, downloads",
    "Bad for: heavy queries, model training, anything CPU-hungry",
], size=14)
# Compute node card
add_rect(s, Inches(6.8), col_y, Inches(6), col_h, LIGHT_GRAY)
add_rect(s, Inches(6.8), col_y, Inches(6), Inches(0.6), MAROON)
add_text(s, Inches(6.8), col_y, Inches(6), Inches(0.6), "COMPUTE NODE",
         font=HEAD_FONT, size=20, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(7.0), col_y + Inches(0.8), Inches(5.6), Inches(4), [
    "Where actual work happens",
    "You request one with idev or sbatch",
    "Your own playground for the duration",
    "Where every query in this lecture will run",
], size=14)
add_text(s, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
         "idev -p gg -m 60 = 'give me a CPU compute node for 60 minutes'",
         font=MONO_FONT, size=14, bold=True, color=MAROON, align=PP_ALIGN.CENTER)
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 9: Filesystem
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "THE TACC FILESYSTEM")
add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
         "Think of it like your desk.",
         font=BODY_FONT, size=20, bold=True, color=MAROON)

table_top = Inches(1.85)
path_w = Inches(1.8)
analogy_w = Inches(3.2)
purpose_w = Inches(7.2)
row_h = Inches(1.2)
hr = Inches(0.6)
x_path, x_analogy, x_purpose = Inches(0.5), Inches(2.4), Inches(5.7)

# Header row
add_rect(s, x_path, table_top, path_w, hr, MAROON)
add_rect(s, x_analogy, table_top, analogy_w, hr, MAROON)
add_rect(s, x_purpose, table_top, purpose_w, hr, MAROON)
add_text(s, x_path, table_top, path_w, hr, "Where",
         font=HEAD_FONT, size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x_analogy, table_top, analogy_w, hr, "Like a...",
         font=HEAD_FONT, size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x_purpose, table_top, purpose_w, hr, "What it holds",
         font=HEAD_FONT, size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

filesystems = [
    ("$HOME", "Drawer",
     "Small (~10 GB), backed up. Login scripts, configs, small files you reuse. Don't keep big data here."),
    ("$WORK", "Desktop",
     "Larger (TBs), for active project files. Today's database lives here. Where most of your work sits."),
    ("$SCRATCH", "Whiteboard",
     "Very large, temporary. Auto-purged every ~10 days since last touch. Use for running jobs, not long-term storage."),
]
y = table_top + hr
for i, (where, analogy, what) in enumerate(filesystems):
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(s, x_path, y, path_w, row_h, bg)
    add_rect(s, x_analogy, y, analogy_w, row_h, bg)
    add_rect(s, x_purpose, y, purpose_w, row_h, bg)
    add_text(s, x_path, y, path_w, row_h, where,
             font=MONO_FONT, size=18, bold=True, color=MAROON,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x_analogy, y, analogy_w, row_h, analogy,
             font=HEAD_FONT, size=20, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x_purpose + Inches(0.15), y, purpose_w - Inches(0.3), row_h, what,
             font=BODY_FONT, size=13, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
    y += row_h

add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
         "Today, your database and all your work go in $WORK.",
         font=BODY_FONT, size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide: You're an analyst
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "TODAY, YOU'RE A DATA ANALYST")
add_text(s, Inches(0.5), Inches(1.5), Inches(12), Inches(0.6),
         "At the NYC Taxi & Limousine Commission.",
         font=BODY_FONT, size=24, bold=True, color=MAROON)
add_text(s, Inches(0.5), Inches(2.4), Inches(12), Inches(0.5),
         "The commission is the regulator. They set the rules every NYC yellow cab follows.",
         font=BODY_FONT, size=18, color=BLACK)
add_text(s, Inches(0.5), Inches(3.5), Inches(12), Inches(0.5),
         "Your boss walks in and asks:",
         font=BODY_FONT, size=18, color=BLACK)
add_rect(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(1.5), LIGHT_GRAY)
add_rect(s, Inches(0.5), Inches(4.3), Inches(0.2), Inches(1.5), GOLD)
add_text(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(1.1),
         '"Should we raise late-night fares? Take a position with evidence by Friday."',
         font=BODY_FONT, size=24, bold=True, color=MAROON, anchor=MSO_ANCHOR.MIDDLE)
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide: The data (and why it needs HPC)
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "THE DATA")
add_text(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
         "All NYC yellow cab trips from 2023.",
         font=BODY_FONT, size=22, bold=True, color=MAROON)

# Three big number callouts
box_y = Inches(2.1)
box_h = Inches(1.7)
box_w = Inches(4.0)
gap = Inches(0.15)
x1 = Inches(0.5)
x2 = x1 + box_w + gap
x3 = x2 + box_w + gap

# Box 1: 38M rows
add_rect(s, x1, box_y, box_w, box_h, MAROON)
add_text(s, x1, box_y + Inches(0.2), box_w, Inches(0.4),
         "ROWS", font=HEAD_FONT, size=14, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x1, box_y + Inches(0.55), box_w, Inches(0.7),
         "38 million", font=HEAD_FONT, size=32, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x1, box_y + Inches(1.25), box_w, Inches(0.4),
         "every yellow cab trip in 2023",
         font=BODY_FONT, size=12, color=COOL_GRAY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Box 2: 6.3 GB DB
add_rect(s, x2, box_y, box_w, box_h, MAROON)
add_text(s, x2, box_y + Inches(0.2), box_w, Inches(0.4),
         "DATABASE", font=HEAD_FONT, size=14, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x2, box_y + Inches(0.55), box_w, Inches(0.7),
         "6.3 GB", font=HEAD_FONT, size=36, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x2, box_y + Inches(1.25), box_w, Inches(0.4),
         "SQLite file on disk",
         font=BODY_FONT, size=12, color=COOL_GRAY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Box 3: 9+ GB to build
add_rect(s, x3, box_y, box_w, box_h, GOLD)
add_text(s, x3, box_y + Inches(0.2), box_w, Inches(0.4),
         "TO BUILD", font=HEAD_FONT, size=14, bold=True, color=MAROON,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x3, box_y + Inches(0.55), box_w, Inches(0.7),
         "~9 GB", font=HEAD_FONT, size=36, bold=True, color=MAROON,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x3, box_y + Inches(1.25), box_w, Inches(0.4),
         "3 GB downloads + 6 GB DB",
         font=BODY_FONT, size=12, color=BLACK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
         "This is why we need HPC.",
         font=HEAD_FONT, size=24, bold=True, color=MAROON, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.5),
         "9 GB of working space and 38M rows are uncomfortable on most laptops, even before you ask any questions.",
         font=BODY_FONT, size=15, color=BLACK, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(5.7), Inches(12), Inches(0.4),
         "Two tables in the database:",
         font=HEAD_FONT, size=14, bold=True, color=MAROON)
add_text(s, Inches(0.5), Inches(6.1), Inches(12), Inches(0.4),
         "trips: pickup_time, dropoff_time, passengers, distance_miles, "
         "pickup_zone_id, dropoff_zone_id, fare, tip, total, payment_type",
         font=MONO_FONT, size=11, color=BLACK)
add_text(s, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
         "zones: zone_id, zone_name (Manhattan neighborhood lookup, 70 rows)",
         font=MONO_FONT, size=11, color=BLACK)

footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 13: Thinking process
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "HOW AN ANALYST ACTUALLY THINKS")
steps = [
    ("1", "Clarify the question", "Whose interest are you weighing? Driver welfare? Revenue? Riders?"),
    ("2", "Argue both sides on paper", "Decide what would make you change your mind BEFORE looking at data."),
    ("3", "Ask: what data could settle each argument?", "Each argument becomes a question. Each question becomes a query."),
    ("4", "Run the queries, interpret the results", "SQL is the tool, not the point."),
    ("5", "Acknowledge what the data CAN'T tell you", "Honest analysts are explicit about the limits."),
]
y = Inches(1.5)
for num, title, desc in steps:
    add_rect(s, Inches(0.5), y, Inches(0.7), Inches(0.7), GOLD)
    add_text(s, Inches(0.5), y, Inches(0.7), Inches(0.7), num,
             font=HEAD_FONT, size=24, bold=True, color=MAROON,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.4), y + Inches(0.05), Inches(11.4), Inches(0.4),
             title, font=HEAD_FONT, size=18, bold=True, color=MAROON)
    add_text(s, Inches(1.4), y + Inches(0.42), Inches(11.4), Inches(0.4),
             desc, font=BODY_FONT, size=14, color=BLACK)
    y += Inches(0.95)
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 14: Switch to terminal
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, BLACK)
add_rect(s, Inches(0.5), Inches(2.0), Inches(0.2), Inches(3.5), GOLD)
add_text(s, Inches(1.0), Inches(2.2), Inches(11.5), Inches(1.0),
         "SWITCH TO TERMINAL", font=HEAD_FONT, size=60, bold=True, color=GOLD)
add_text(s, Inches(1.0), Inches(3.5), Inches(11.5), Inches(0.6),
         "Live demo: 6 queries, 5 phases, one investigation.",
         font=BODY_FONT, size=22, color=WHITE)
add_text(s, Inches(1.0), Inches(4.5), Inches(11.5), Inches(0.5),
         "Follow along: ssh username@vista.tacc.utexas.edu",
         font=MONO_FONT, size=18, color=COOL_GRAY)

# =============================================================================
# Slide 15: Phase 1 recap
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "PHASE 1: MEET THE DATA")
add_text(s, Inches(0.5), Inches(1.4), Inches(12), Inches(0.5),
         "Question: what are we actually working with?",
         font=BODY_FONT, size=20, bold=True, color=MAROON)
# query box
add_rect(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.0), LIGHT_GRAY)
add_text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(0.4),
         "SELECT COUNT(*) FROM trips;",
         font=MONO_FONT, size=18, bold=True, color=MAROON)
add_text(s, Inches(0.7), Inches(2.7), Inches(12), Inches(0.4),
         "SELECT * FROM trips LIMIT 5;",
         font=MONO_FONT, size=18, bold=True, color=MAROON)
add_text(s, Inches(0.5), Inches(3.6), Inches(12), Inches(0.5),
         "What we found:",
         font=HEAD_FONT, size=18, bold=True, color=MAROON)
add_bullets(s, Inches(0.5), Inches(4.1), Inches(12), Inches(2.5), [
    "38,310,226 trips in 2023.",
    "pickup_time is stored as TEXT, not a real timestamp. We'll do string surgery later.",
    "payment_type is coded: 1=Credit, 2=Cash, 3=No charge, 4=Dispute.",
], size=18)
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 16: Phase 2 recap
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "PHASE 2: WHEN DO TRIPS HAPPEN?")
add_text(s, Inches(0.5), Inches(1.4), Inches(12), Inches(0.5),
         "Question: are late-night trips even a meaningful slice of business?",
         font=BODY_FONT, size=20, bold=True, color=MAROON)
add_text(s, Inches(0.5), Inches(2.2), Inches(12), Inches(0.5),
         "What we found:",
         font=HEAD_FONT, size=18, bold=True, color=MAROON)
add_bullets(s, Inches(0.5), Inches(2.7), Inches(12), Inches(3), [
    "Peak hour: 6 PM with 2.7 million trips. Trough: 4 AM with 217,000.",
    "Late-night (10 PM to 4:59 AM) = 6.4M trips, or 16.7% of total business.",
    "Not trivial, not dominant. Big enough to matter for a fare-change decision.",
], size=18)
# big number callout
add_rect(s, Inches(8.5), Inches(5.5), Inches(4), Inches(1.4), MAROON)
add_text(s, Inches(8.5), Inches(5.55), Inches(4), Inches(0.5), "LATE-NIGHT SHARE",
         font=HEAD_FONT, size=14, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(8.5), Inches(5.95), Inches(4), Inches(0.9), "16.7%",
         font=HEAD_FONT, size=44, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 17: Phase 3 recap
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "PHASE 3: WHERE DO LATE-NIGHT TRIPS HAPPEN?")
add_text(s, Inches(0.5), Inches(1.4), Inches(12), Inches(0.5),
         "Question: who's actually riding at night?",
         font=BODY_FONT, size=20, bold=True, color=MAROON)
zones = [
    ("1", "JFK Airport", "416,574"),
    ("2", "East Village", "402,029"),
    ("3", "West Village", "339,860"),
    ("4", "Clinton East", "281,816"),
    ("5", "Sutton Place / Turtle Bay South", "254,150"),
    ("6", "Lower East Side", "236,074"),
    ("7", "Greenwich Village South", "230,367"),
    ("8", "Midtown Center", "206,353"),
    ("9", "Lincoln Square East", "199,624"),
    ("10", "Penn Station / Madison Sq West", "191,701"),
]
y = Inches(2.2)
for rank, name, count in zones:
    add_text(s, Inches(0.5), y, Inches(0.5), Inches(0.3), rank,
             font=HEAD_FONT, size=14, bold=True, color=MAROON, align=PP_ALIGN.RIGHT)
    add_text(s, Inches(1.1), y, Inches(8), Inches(0.3), name,
             font=BODY_FONT, size=14, color=BLACK)
    add_text(s, Inches(8.5), y, Inches(2), Inches(0.3), count,
             font=MONO_FONT, size=14, color=GOLD, align=PP_ALIGN.RIGHT)
    y += Inches(0.32)
add_text(s, Inches(0.5), Inches(6.3), Inches(12.5), Inches(0.5),
         "Mixed picture: airport (red-eyes), entertainment districts, transit hubs. NOT just bar crowd.",
         font=BODY_FONT, size=16, bold=True, color=MAROON)
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 18: Phase 4 recap (THE SURPRISE)
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "PHASE 4: HOW ARE LATE-NIGHT TRIPS DIFFERENT?")
add_text(s, Inches(0.5), Inches(1.4), Inches(12), Inches(0.5),
         "Question: in fare, distance, and tipping, is late-night a different beast?",
         font=BODY_FONT, size=20, bold=True, color=MAROON)
# comparison table
table_top = Inches(2.2)
hr = Inches(0.7)
cw1 = Inches(3.5)
cw2 = Inches(3.0)
cw3 = Inches(3.0)
cw4 = Inches(2.8)
x1, x2, x3, x4 = Inches(0.5), Inches(4.0), Inches(7.0), Inches(10.0)

add_rect(s, x1, table_top, cw1, hr, MAROON)
add_rect(s, x2, table_top, cw2, hr, MAROON)
add_rect(s, x3, table_top, cw3, hr, MAROON)
add_rect(s, x4, table_top, cw4, hr, MAROON)
add_text(s, x1, table_top, cw1, hr, "Metric", font=HEAD_FONT, size=18, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x2, table_top, cw2, hr, "Daytime", font=HEAD_FONT, size=18, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x3, table_top, cw3, hr, "Late-Night", font=HEAD_FONT, size=18, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x4, table_top, cw4, hr, "Difference", font=HEAD_FONT, size=18, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

rows = [
    ("Avg fare", "$19.51", "$19.60", "+$0.09 (0.5%)"),
    ("Avg distance", "4.04 mi", "4.32 mi", "+0.28 mi (7%)"),
    ("Avg tip", "$3.52", "$3.51", "$0.01"),
]
y = table_top + hr
for i, (label, day, night, diff) in enumerate(rows):
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(s, x1, y, cw1, hr, bg)
    add_rect(s, x2, y, cw2, hr, bg)
    add_rect(s, x3, y, cw3, hr, bg)
    add_rect(s, x4, y, cw4, hr, bg)
    add_text(s, x1 + Inches(0.1), y, cw1, hr, label,
             font=HEAD_FONT, size=14, bold=True, color=MAROON, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x2, y, cw2, hr, day, font=BODY_FONT, size=14, color=BLACK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x3, y, cw3, hr, night, font=BODY_FONT, size=14, color=BLACK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x4, y, cw4, hr, diff, font=BODY_FONT, size=14, color=GOLD,
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += hr

add_text(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.6),
         "The surprise: late-night and daytime trips are basically identical.",
         font=HEAD_FONT, size=22, bold=True, color=MAROON, align=PP_ALIGN.CENTER)
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 19: Saving and handing off your work
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "SAVE YOUR WORK, HAND IT OFF")
add_text(s, Inches(0.5), Inches(1.4), Inches(12), Inches(0.5),
         "An analysis is only useful if someone else can pick it up.",
         font=BODY_FONT, size=20, bold=True, color=MAROON)
add_text(s, Inches(0.5), Inches(2.2), Inches(12), Inches(0.5),
         "In SQLite, save each query's output to a CSV:",
         font=BODY_FONT, size=16, color=BLACK)
add_rect(s, Inches(0.5), Inches(2.8), Inches(12.3), Inches(1.6), LIGHT_GRAY)
code_lines = [
    ".headers on",
    ".mode csv",
    ".once query_1.csv",
    "SELECT COUNT(*) FROM trips;",
]
y = Inches(2.9)
for line in code_lines:
    add_text(s, Inches(0.7), y, Inches(12), Inches(0.35), line,
             font=MONO_FONT, size=16, bold=True, color=MAROON)
    y += Inches(0.32)

add_text(s, Inches(0.5), Inches(4.7), Inches(12), Inches(0.5),
         "Then view your saved files from the shell:",
         font=BODY_FONT, size=16, color=BLACK)
add_rect(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.6), LIGHT_GRAY)
shell_lines = [
    "ls *.csv                           # see what we saved",
    "less query_2.csv                   # paginated view",
    "column -t -s, query_5.csv | less   # aligned columns",
]
y = Inches(5.4)
for line in shell_lines:
    add_text(s, Inches(0.7), y, Inches(12), Inches(0.35), line,
             font=MONO_FONT, size=14, bold=True, color=MAROON)
    y += Inches(0.32)

footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 20: Recommendation
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "THE RECOMMENDATION")
add_rect(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.0), MAROON)
add_rect(s, Inches(0.5), Inches(1.5), Inches(0.2), Inches(2.0), GOLD)
add_text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.5),
         "Based on this data alone:",
         font=BODY_FONT, size=18, color=GOLD)
add_text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.2),
         "Don't raise late-night fares.",
         font=HEAD_FONT, size=42, bold=True, color=WHITE)
add_text(s, Inches(0.5), Inches(3.9), Inches(12), Inches(0.5),
         "The case the numbers make:",
         font=HEAD_FONT, size=18, bold=True, color=MAROON)
add_bullets(s, Inches(0.5), Inches(4.4), Inches(12), Inches(2.3), [
    "16.7% of trips are late-night, but they're nearly indistinguishable from daytime in fare, distance, and tipping.",
    "Riders aren't paying more. Drivers aren't earning visibly more. Tips are flat.",
    "Whatever case there is for raising fares, it cannot be made on these averages.",
], size=16)
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 21: Limitations
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "WHAT THE DATA CANNOT TELL US")
add_text(s, Inches(0.5), Inches(1.4), Inches(12), Inches(0.5),
         "An honest analyst is explicit about the limits.",
         font=BODY_FONT, size=20, bold=True, color=MAROON)
add_bullets(s, Inches(0.5), Inches(2.2), Inches(12), Inches(4.5), [
    "Driver welfare: no driver IDs, no hours worked, no take-home pay.",
    "Cash tips: payment_type 2 records $0 tips, biasing the tip analysis.",
    "Rider demographics or income.",
    "Demand elasticity: we only see current prices, not how riders would respond to a change.",
    "Alternative transit / Uber-Lyft competition.",
    "Outer boroughs: zones table is mostly Manhattan + airports.",
], size=18)
add_text(s, Inches(0.5), Inches(6.6), Inches(12), Inches(0.5),
         "The recommendation should name these gaps and suggest what data would strengthen the case.",
         font=BODY_FONT, size=14, color=COOL_GRAY)
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 22: Honest analyst's job
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, MAROON)
add_rect(s, Inches(0.5), Inches(2.5), Inches(0.2), Inches(3), GOLD)
add_text(s, Inches(1.0), Inches(2.5), Inches(11.5), Inches(0.5),
         "THE JOB", font=HEAD_FONT, size=20, bold=True, color=GOLD)
add_text(s, Inches(1.0), Inches(3.0), Inches(11.5), Inches(2.5),
         "The data sometimes refuses to support the obvious answer.",
         font=HEAD_FONT, size=40, bold=True, color=WHITE)
add_text(s, Inches(1.0), Inches(5.5), Inches(11.5), Inches(0.5),
         "Your job is to be honest about that.",
         font=BODY_FONT, size=22, color=GOLD)

# =============================================================================
# Slide: Homework
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "HOMEWORK")
add_text(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
         "Same dataset. New question from the commission:",
         font=BODY_FONT, size=18, color=BLACK)

# The question in a callout box
add_rect(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.95), MAROON)
add_rect(s, Inches(0.5), Inches(1.95), Inches(0.2), Inches(0.95), GOLD)
add_text(s, Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.95),
         "Should airport fares be priced differently from non-airport fares?",
         font=HEAD_FONT, size=22, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.5), Inches(3.1), Inches(12), Inches(0.4),
         "Airport zones: pickup_zone_id IN (1, 132, 138) — Newark, JFK, LaGuardia",
         font=MONO_FONT, size=13, color=COOL_GRAY)

add_text(s, Inches(0.5), Inches(3.7), Inches(12), Inches(0.4),
         "Sub-questions to guide your queries:",
         font=HEAD_FONT, size=16, bold=True, color=MAROON)
add_bullets(s, Inches(0.5), Inches(4.2), Inches(12), Inches(2), [
    "How big a slice of the business are airport trips?",
    "How are airport trips different in fare, distance, tip?",
    "Are airport trips concentrated at certain hours of the day?",
    "Are airport riders more likely to pay by credit card vs cash?",
], size=14)

add_text(s, Inches(0.5), Inches(6.1), Inches(12), Inches(0.4),
         "Submit:",
         font=HEAD_FONT, size=14, bold=True, color=MAROON)
add_text(s, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
         "CSV outputs  +  recommendation (1-2 paragraphs)  +  reflection (~150 words)",
         font=BODY_FONT, size=13, color=BLACK)
footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide: What else can you do with HPC
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
header_bar(s, "WHAT ELSE CAN YOU DO WITH HPC?")
add_text(s, Inches(0.5), Inches(1.4), Inches(12), Inches(0.5),
         "Today was SQL on a 6 GB database. The system can do much more.",
         font=BODY_FONT, size=18, color=BLACK)

use_cases = [
    ("Train AI/ML models",
     "Hardware built for it. Days on a laptop become hours on HPC."),
    ("Run hundreds of experiments at once",
     "Parallel jobs, not waiting for one to finish before starting the next."),
    ("Work with large protected datasets",
     "Hosted in a secure environment. Analyze without copying."),
    ("Generate synthetic data at scale",
     "Synthetic images, text, simulations in batch jobs."),
    ("Render large-scale visualizations",
     "Complex figures and animations that would crash your laptop."),
]
y = Inches(2.3)
for title, desc in use_cases:
    add_rect(s, Inches(0.5), y, Inches(0.15), Inches(0.7), GOLD)
    add_text(s, Inches(0.85), y + Inches(0.0), Inches(11.5), Inches(0.35),
             title, font=HEAD_FONT, size=18, bold=True, color=MAROON)
    add_text(s, Inches(0.85), y + Inches(0.38), Inches(11.5), Inches(0.35),
             desc, font=BODY_FONT, size=14, color=BLACK)
    y += Inches(0.85)

footer(s, slide_num, TOTAL_SLIDES)

# =============================================================================
# Slide 26: Q&A
# =============================================================================
slide_num += 1
s = prs.slides.add_slide(BLANK)
add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, MAROON)
add_rect(s, Inches(0.5), Inches(3.0), Inches(2), Inches(0.1), GOLD)
add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.0),
         "Q&A", font=HEAD_FONT, size=120, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.6),
         "Thank you.", font=BODY_FONT, size=36, color=GOLD,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
         "Ashley Scruse, Ph.D.  |  Morehouse Supercomputing Facility",
         font=BODY_FONT, size=14, color=COOL_GRAY, align=PP_ALIGN.CENTER)

# Save
import os
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "sql_on_hpc.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {slide_num}")
